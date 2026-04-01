"""Generate The Ppitch Awards 2026 Presentation — Hospital Claim AI
HIGH PERFORMANCE HEALTHCARE | รพ.พญาไทศรีราชา | BDMS G5
FONT RULES: Title 40pt+, Body 20pt+, Stats 64pt+, Min 18pt
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# ═══════════════════════════════════════
# Colors
# ═══════════════════════════════════════
NAVY = RGBColor(0x0A, 0x1F, 0x3B)
TEAL = RGBColor(0x00, 0x8B, 0x8B)
GOLD = RGBColor(0xD4, 0xA5, 0x37)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xF5, 0xF5, 0xF0)
DARK = RGBColor(0x1A, 0x1A, 0x2E)
GRAY = RGBColor(0x8B, 0x8B, 0x8B)
RED = RGBColor(0xE8, 0x3E, 0x3E)
GREEN = RGBColor(0x00, 0xA8, 0x6B)
BLUE = RGBColor(0x2E, 0x86, 0xC1)
ORANGE = RGBColor(0xE6, 0x7E, 0x22)

W = Inches(13.333)
H = Inches(7.5)


def bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def rect(slide, l, t, w, h, color):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s


def rrect(slide, l, t, w, h, color):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s


def txt(slide, l, t, w, h, text, sz=20, bold=False, color=DARK, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tb.text_frame.word_wrap = True
    p = tb.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = "Arial Black" if bold and sz >= 36 else "Arial"
    p.alignment = align
    return tb


def circle(slide, l, t, sz, color, text, txt_sz=24):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, l, t, Inches(sz), Inches(sz))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    p = s.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(txt_sz)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    return s


def stat(slide, l, t, w, h, number, label, bg_color):
    rrect(slide, l, t, w, h, bg_color)
    txt(slide, l, t + Inches(0.15), w, Inches(1.2), number, 64, True, WHITE, PP_ALIGN.CENTER)
    txt(slide, l, t + h - Inches(0.55), w, Inches(0.5), label, 20, True, WHITE, PP_ALIGN.CENTER)


def title_bar(slide, text, color=NAVY):
    rect(slide, Inches(0), Inches(0), W, Inches(1.2), color)
    txt(slide, Inches(0.8), Inches(0.2), Inches(11.5), Inches(0.8), text, 40, True, WHITE)


def generate():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    # ═══════════════════════════════════
    # S1: COVER
    # ═══════════════════════════════════
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, NAVY)
    rect(s, Inches(0), Inches(0), W, Inches(0.1), GOLD)

    rrect(s, Inches(3.8), Inches(1.0), Inches(5.7), Inches(0.7), GOLD)
    txt(s, Inches(3.8), Inches(1.05), Inches(5.7), Inches(0.6),
        "HIGH PERFORMANCE HEALTHCARE", 24, True, NAVY, PP_ALIGN.CENTER)

    txt(s, Inches(0.5), Inches(2.2), Inches(12.3), Inches(1.2),
        "Hospital Claim AI", 72, True, WHITE, PP_ALIGN.CENTER)

    txt(s, Inches(1), Inches(3.6), Inches(11.3), Inches(0.7),
        "ลด Deny เพิ่มรายได้ รองรับทุกกองทุน ทุกแผนก", 28, False, GOLD, PP_ALIGN.CENTER)

    txt(s, Inches(1), Inches(4.6), Inches(11.3), Inches(0.5),
        "โรงพยาบาลพญาไทศรีราชา  |  รหัส 11855  |  ประกันสังคม + สปสช. + ประกันเอกชน", 22, False, GRAY, PP_ALIGN.CENTER)

    txt(s, Inches(1), Inches(5.5), Inches(11.3), Inches(0.5),
        "ทีม: แผนกศัลยกรรมหัวใจและทรวงอก", 22, False, WHITE, PP_ALIGN.CENTER)
    txt(s, Inches(1), Inches(6.0), Inches(11.3), Inches(0.5),
        "นพ.เกษมสันต์ เกษมวงศ์  |  รหัสพนักงาน 3241", 20, False, GRAY, PP_ALIGN.CENTER)

    rect(s, Inches(0), Inches(7.0), W, Inches(0.5), TEAL)
    txt(s, Inches(0), Inches(7.05), W, Inches(0.4),
        "The Ppitch Awards 2026  |  G5 Passion to Innovate  |  BDMS", 18, True, WHITE, PP_ALIGN.CENTER)

    # ═══════════════════════════════════
    # S2: PAIN POINTS
    # ═══════════════════════════════════
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, LIGHT)
    title_bar(s, "Pain Points — ปัญหาที่ทำให้ รพ. สูญเสียรายได้")

    # Left: 4 problems
    probs = [
        ("1", "ICD Coding ไม่ตรง Procedure", RED),
        ("2", "เอกสารอุปกรณ์ไม่ครบ (HC09/HC13)", RED),
        ("3", "ตรวจ Manual 15-30 นาที/เคส", ORANGE),
        ("4", "อุทธรณ์ไม่ทัน ไม่เป็นระบบ", ORANGE),
    ]
    for i, (num, text, color) in enumerate(probs):
        y = Inches(1.6) + Inches(i * 1.35)
        circle(s, Inches(0.8), y, 0.65, color, num, 28)
        txt(s, Inches(1.7), y + Inches(0.05), Inches(5.5), Inches(0.6), text, 24, True, DARK)

    # Right: Impact stats on dark bg
    rect(s, Inches(7.5), Inches(1.4), Inches(5.3), Inches(5.8), NAVY)
    txt(s, Inches(7.5), Inches(1.6), Inches(5.3), Inches(0.5),
        "ผลกระทบ", 28, True, GOLD, PP_ALIGN.CENTER)

    impacts = [
        ("~10%", "Deny Rate"),
        ("4.8M", "บาท/ปี สูญเสีย"),
        ("30 นาที", "ตรวจ 1 เคส"),
        ("2 ชม.", "ร่าง Appeal 1 เคส"),
    ]
    for i, (num, label) in enumerate(impacts):
        y = Inches(2.4) + Inches(i * 1.2)
        txt(s, Inches(7.8), y, Inches(4.8), Inches(0.7), num, 48, True, GOLD, PP_ALIGN.CENTER)
        txt(s, Inches(7.8), y + Inches(0.6), Inches(4.8), Inches(0.4), label, 20, False, WHITE, PP_ALIGN.CENTER)

    # ═══════════════════════════════════
    # S3: SOLUTION — 8 Checkpoints
    # ═══════════════════════════════════
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, LIGHT)
    title_bar(s, "Solution — AI ตรวจ 8 จุดก่อนส่งเบิก", TEAL)

    cps = [
        ("01", "ข้อมูลพื้นฐาน", BLUE),
        ("02", "Dx-Proc Match", BLUE),
        ("03", "เอกสารอุปกรณ์", BLUE),
        ("04", "FDH 16-File", BLUE),
        ("05", "Timing + Authen", ORANGE),
        ("06", "CC/MCC Optimize", ORANGE),
        ("07", "DRG Verify", ORANGE),
        ("08", "Drug/Lab Catalog", ORANGE),
    ]
    for i, (num, name, color) in enumerate(cps):
        col = i % 4
        row = i // 4
        x = Inches(0.5) + Inches(col * 3.2)
        y = Inches(1.5) + Inches(row * 2.9)

        rrect(s, x, y, Inches(2.9), Inches(2.5), WHITE)
        circle(s, x + Inches(0.95), y + Inches(0.2), 0.8, color, num, 30)
        txt(s, x + Inches(0.1), y + Inches(1.2), Inches(2.7), Inches(0.8),
            name, 24, True, DARK, PP_ALIGN.CENTER)

    # ═══════════════════════════════════
    # S4: 6 FEATURES
    # ═══════════════════════════════════
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, LIGHT)
    title_bar(s, "6 ฟีเจอร์ — คำสั่งเดียว AI จัดการ")

    feats = [
        ("/claim check", "ตรวจก่อนส่ง", "Score 0-100%", BLUE),
        ("/claim predict", "ทำนาย Deny", "10 Risk Factors", TEAL),
        ("/claim code", "Smart Coder", "Notes → ICD auto", GREEN),
        ("/claim deny", "วิเคราะห์ Deny", "Root Cause + Fix", RED),
        ("/claim appeal", "ร่าง Appeal", "DOCX พร้อมพิมพ์", ORANGE),
        ("/claim batch", "Batch Optimize", "CSV ตรวจทั้งล็อต", NAVY),
    ]
    for i, (cmd, title, desc, color) in enumerate(feats):
        col = i % 3
        row = i // 3
        x = Inches(0.5) + Inches(col * 4.2)
        y = Inches(1.5) + Inches(row * 3.0)

        rrect(s, x, y, Inches(3.8), Inches(2.6), WHITE)
        rrect(s, x + Inches(0.2), y + Inches(0.2), Inches(3.0), Inches(0.5), color)
        txt(s, x + Inches(0.2), y + Inches(0.22), Inches(3.0), Inches(0.5),
            cmd, 20, True, WHITE, PP_ALIGN.CENTER)
        txt(s, x + Inches(0.2), y + Inches(0.9), Inches(3.4), Inches(0.5),
            title, 28, True, DARK)
        txt(s, x + Inches(0.2), y + Inches(1.5), Inches(3.4), Inches(0.5),
            desc, 22, False, GRAY)

    # ═══════════════════════════════════
    # S5: BENCHMARK
    # ═══════════════════════════════════
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, LIGHT)
    title_bar(s, "Benchmark — ก่อน vs หลังใช้ AI")

    # Before
    rect(s, Inches(0.5), Inches(1.5), Inches(5.9), Inches(5.5), WHITE)
    rrect(s, Inches(1.5), Inches(1.6), Inches(3.8), Inches(0.6), RED)
    txt(s, Inches(1.5), Inches(1.63), Inches(3.8), Inches(0.55),
        "BEFORE", 24, True, WHITE, PP_ALIGN.CENTER)

    befores = [
        "ตรวจ 15-30 นาที/เคส",
        "เปิดหลายระบบ ตรวจด้วยตา",
        "Deny Rate ~10%",
        "ร่าง Appeal 1-2 ชั่วโมง",
        "สูญเสีย 4.8M บาท/ปี",
    ]
    for i, item in enumerate(befores):
        txt(s, Inches(1.0), Inches(2.5) + Inches(i * 0.8), Inches(5.0), Inches(0.6),
            f"  {item}", 22, False, RED)

    # After
    rect(s, Inches(6.9), Inches(1.5), Inches(5.9), Inches(5.5), WHITE)
    rrect(s, Inches(7.9), Inches(1.6), Inches(3.8), Inches(0.6), GREEN)
    txt(s, Inches(7.9), Inches(1.63), Inches(3.8), Inches(0.55),
        "AFTER — AI", 24, True, WHITE, PP_ALIGN.CENTER)

    afters = [
        "ตรวจ < 5 นาที/เคส",
        "AI ตรวจ 8 checkpoints อัตโนมัติ",
        "Deny Rate < 3%",
        "ร่าง Appeal 5 นาที + DOCX",
        "เพิ่มรายได้ 2.5M+ บาท/ปี",
    ]
    for i, item in enumerate(afters):
        txt(s, Inches(7.4), Inches(2.5) + Inches(i * 0.8), Inches(5.0), Inches(0.6),
            f"  {item}", 22, False, GREEN)

    # ═══════════════════════════════════
    # S6: CASE STUDY
    # ═══════════════════════════════════
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, LIGHT)
    title_bar(s, "Case Study — เคสจริง STEMI + PCI", TEAL)

    # Patient card
    rect(s, Inches(0.5), Inches(1.5), Inches(6.2), Inches(5.5), WHITE)
    txt(s, Inches(0.8), Inches(1.6), Inches(5.6), Inches(0.5),
        "AN 69-03556  |  STEMI Anterior", 24, True, DARK)

    lines = [
        "PDx: I21.0 Acute MI anterior wall",
        "Proc: 36.06, 37.22, 88.56 (PCI+Cath+Angio)",
        "DRG: 05290 | RW: 8.6544",
        "Charge: 71,322 บาท",
    ]
    for i, line in enumerate(lines):
        txt(s, Inches(0.8), Inches(2.3) + Inches(i * 0.7), Inches(5.6), Inches(0.5),
            line, 22, False, DARK)

    rrect(s, Inches(0.8), Inches(5.3), Inches(5.6), Inches(0.6), RED)
    txt(s, Inches(0.8), Inches(5.33), Inches(5.6), Inches(0.55),
        "Deny: HC09 + HC13 + Drug Code ไม่ตรง", 20, True, WHITE, PP_ALIGN.CENTER)

    # AI result
    rect(s, Inches(7.2), Inches(1.5), Inches(5.6), Inches(5.5), NAVY)
    txt(s, Inches(7.5), Inches(1.6), Inches(5.0), Inches(0.5),
        "AI ตรวจพบ + แก้ไขได้:", 24, True, GOLD)

    fixes = [
        "ADP: TYPE=5 + Serial ตรง GPO",
        "DRU: TMT code ต้อง map ใหม่",
        "แนะนำ: แก้แล้วส่งใหม่",
    ]
    for i, fix in enumerate(fixes):
        txt(s, Inches(7.5), Inches(2.4) + Inches(i * 0.7), Inches(5.0), Inches(0.5),
            fix, 22, False, WHITE)

    # Score cards
    stat(s, Inches(7.4), Inches(4.4), Inches(2.4), Inches(2.2), "65%", "ก่อนแก้", RED)
    stat(s, Inches(10.2), Inches(4.4), Inches(2.4), Inches(2.2), "12%", "หลังแก้", GREEN)

    # ═══════════════════════════════════
    # S7: RESULTS — Big Numbers
    # ═══════════════════════════════════
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, NAVY)

    txt(s, Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.8),
        "Results — ผลลัพธ์ที่วัดได้", 44, True, GOLD, PP_ALIGN.CENTER)

    results = [
        (">90%", "เบิกผ่านเพิ่ม", GREEN),
        ("+15%", "DRG Weight", TEAL),
        ("-83%", "ลดเวลาตรวจ", BLUE),
        ("2.5M+", "บาท/ปี เพิ่ม", GOLD),
    ]
    for i, (num, label, color) in enumerate(results):
        x = Inches(0.4) + Inches(i * 3.2)
        stat(s, x, Inches(1.5), Inches(2.9), Inches(2.5), num, label, color)

    # Details
    rows = [
        ("เวลาตรวจ claim", "15-30 นาที → 5 นาที"),
        ("เวลาร่าง Appeal", "1-2 ชม. → 5 นาที"),
        ("รองรับ 3 กองทุน", "ประกันสังคม + สปสช. + ประกันเอกชน"),
        ("AI Skills", "20 skills ครอบคลุม 9 แผนก"),
    ]
    for i, (item, value) in enumerate(rows):
        y = Inches(4.5) + Inches(i * 0.7)
        txt(s, Inches(0.8), y, Inches(4), Inches(0.5), item, 22, True, WHITE)
        txt(s, Inches(5.2), y, Inches(7.5), Inches(0.5), value, 22, False, GOLD)

    # ═══════════════════════════════════
    # S8: MARKET + SCALING
    # ═══════════════════════════════════
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, LIGHT)
    title_bar(s, "ความร่วมมือ & แผนขยายผล")

    # Collab
    rect(s, Inches(0.5), Inches(1.5), Inches(6), Inches(2.5), WHITE)
    txt(s, Inches(0.8), Inches(1.6), Inches(5.4), Inches(0.5),
        "ร่วมพัฒนากับ:", 24, True, TEAL)
    collabs = ["Cath Lab team — ข้อมูล clinical", "IT — HIS (SSB) + infra", "UR/Coder — ทดสอบ + feedback", "สปสช. — Thai DRG + FDH spec"]
    for i, c in enumerate(collabs):
        txt(s, Inches(0.8), Inches(2.2) + Inches(i * 0.45), Inches(5.4), Inches(0.4), c, 20, False, DARK)

    # Market
    rect(s, Inches(6.9), Inches(1.5), Inches(6), Inches(2.5), WHITE)
    txt(s, Inches(7.2), Inches(1.6), Inches(5.4), Inches(0.5),
        "รองรับ 3 กองทุน:", 24, True, TEAL)
    markets = ["ประกันสังคม (SSO) — เคสหลักของ รพ.เอกชน", "สปสช. (UC) — UCEP ฉุกเฉิน", "ประกันเอกชน — DRG logic เดียวกัน"]
    for i, m in enumerate(markets):
        txt(s, Inches(7.2), Inches(2.2) + Inches(i * 0.45), Inches(5.4), Inches(0.4), m, 20, False, DARK)

    # Roadmap phases
    phases = [
        ("Phase 1\nQ3-Q4/69", "ขยายทุก 9 แผนก\nภายใน รพ.", BLUE),
        ("Phase 2\nปี 70", "+ ประกันสังคม\n+ BDMS 3-5 รพ.", TEAL),
        ("Phase 3\nปี 71", "SaaS Platform\n+ ประกันเอกชน", ORANGE),
        ("Phase 4\nปี 72-73", "เปิด รพ. นอกเครือ\nNLP Thai", GOLD),
    ]
    for i, (phase, desc, color) in enumerate(phases):
        x = Inches(0.5) + Inches(i * 3.2)
        rrect(s, x, Inches(4.5), Inches(2.9), Inches(2.5), color)
        txt(s, x, Inches(4.6), Inches(2.9), Inches(0.8), phase, 20, True, WHITE, PP_ALIGN.CENTER)
        txt(s, x, Inches(5.5), Inches(2.9), Inches(1.0), desc, 20, False, WHITE, PP_ALIGN.CENTER)

    # ═══════════════════════════════════
    # S9: AI COMMAND CENTER
    # ═══════════════════════════════════
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, NAVY)

    txt(s, Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.7),
        "Vision — AI Command Center", 44, True, GOLD, PP_ALIGN.CENTER)
    txt(s, Inches(0.5), Inches(0.9), Inches(12.3), Inches(0.5),
        "จาก Claim AI → สมองกลางบริหาร รพ. ทั้งระบบ", 24, False, WHITE, PP_ALIGN.CENTER)

    modules = [
        ("Claim AI", "ลด Deny เพิ่มรายได้", GREEN, "DONE"),
        ("Bed AI", "จัดเตียง ลด Wait", BLUE, "NEXT"),
        ("Staff AI", "จัดเวร ลด OT", TEAL, "PLAN"),
        ("Supply AI", "คลังยา/วัสดุ", ORANGE, "PLAN"),
        ("Finance AI", "งบ + Cashflow", GOLD, "PLAN"),
        ("Quality AI", "JCI/HA + KPI", RED, "PLAN"),
    ]
    for i, (name, desc, color, status) in enumerate(modules):
        col = i % 3
        row = i // 3
        x = Inches(0.8) + Inches(col * 4.1)
        y = Inches(1.7) + Inches(row * 2.5)

        rrect(s, x, y, Inches(3.5), Inches(2.1), color)

        # Status badge
        bc = GREEN if status == "DONE" else (BLUE if status == "NEXT" else GRAY)
        rrect(s, x + Inches(2.3), y + Inches(0.1), Inches(1.0), Inches(0.4), bc)
        txt(s, x + Inches(2.3), y + Inches(0.12), Inches(1.0), Inches(0.35),
            status, 16, True, WHITE, PP_ALIGN.CENTER)

        txt(s, x + Inches(0.2), y + Inches(0.15), Inches(2.0), Inches(0.5),
            name, 28, True, WHITE)
        txt(s, x + Inches(0.2), y + Inches(0.7), Inches(3.1), Inches(0.5),
            desc, 22, False, WHITE)

    # CEO Dashboard
    rrect(s, Inches(3), Inches(6.5), Inches(7.3), Inches(0.7), GOLD)
    txt(s, Inches(3), Inches(6.53), Inches(7.3), Inches(0.65),
        "CEO Dashboard — เห็นทุกอย่าง ตัดสินใจทันที", 24, True, NAVY, PP_ALIGN.CENTER)

    # ═══════════════════════════════════
    # S10: CLOSING
    # ═══════════════════════════════════
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, NAVY)
    rect(s, Inches(0), Inches(0), W, Inches(0.1), GOLD)

    txt(s, Inches(0.5), Inches(1.0), Inches(12.3), Inches(1.2),
        "Hospital Claim AI", 72, True, WHITE, PP_ALIGN.CENTER)

    txt(s, Inches(1), Inches(2.5), Inches(11.3), Inches(0.7),
        "ลด Deny เพิ่มรายได้ รองรับทุกกองทุน ทุกแผนก", 28, False, GOLD, PP_ALIGN.CENTER)

    # Key stats
    closing = [
        ("20", "AI Skills", TEAL),
        ("3", "กองทุน", BLUE),
        ("9", "แผนก", GREEN),
        ("<5 min", "ต่อเคส", ORANGE),
    ]
    for i, (num, label, color) in enumerate(closing):
        x = Inches(0.8) + Inches(i * 3.1)
        stat(s, x, Inches(3.5), Inches(2.7), Inches(2.0), num, label, color)

    txt(s, Inches(0.5), Inches(5.9), Inches(12.3), Inches(0.5),
        "1 แผนก → ทุกกองทุน → ทั้ง BDMS → ทั้งประเทศ", 28, True, GOLD, PP_ALIGN.CENTER)

    txt(s, Inches(1), Inches(6.5), Inches(11.3), Inches(0.4),
        "นพ.เกษมสันต์ เกษมวงศ์  |  kasemson_kas@phyathai.com  |  083-574-7685", 18, False, GRAY, PP_ALIGN.CENTER)

    rect(s, Inches(0), Inches(7.0), W, Inches(0.5), GOLD)
    txt(s, Inches(0), Inches(7.05), W, Inches(0.4),
        "The Ppitch Awards 2026  |  G5 Passion to Innovate  |  BDMS", 18, True, NAVY, PP_ALIGN.CENTER)

    # Save
    out = "docs/Ppitch_Awards_2026_Hospital_Claim_AI.pptx"
    prs.save(out)
    print(f"สร้างเสร็จ: {out}")
    print(f"ขนาด: {os.path.getsize(out):,} bytes | Slides: {len(prs.slides)}")


if __name__ == "__main__":
    generate()
