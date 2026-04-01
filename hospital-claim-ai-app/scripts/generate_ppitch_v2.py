"""The Ppitch Awards 2026 — Hospital Claim AI v2
BETTER: Storytelling flow, bigger fonts, cleaner layout, 3 กองทุน
Scoring focus: Results 45% > Collaboration 30% > Concept 20% > Story 5%
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# Colors — Dark Premium
BG = RGBColor(0x0D, 0x1B, 0x2A)        # Deep dark blue
CARD = RGBColor(0x15, 0x2A, 0x3E)      # Card bg
GOLD = RGBColor(0xF0, 0xB9, 0x40)      # Gold accent
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF8, 0xF9, 0xFA)  # Light slides
DARK = RGBColor(0x1A, 0x1A, 0x2E)
GRAY = RGBColor(0x9B, 0x9B, 0x9B)
RED = RGBColor(0xFF, 0x4D, 0x4D)
GREEN = RGBColor(0x00, 0xC9, 0x7B)
BLUE = RGBColor(0x3A, 0x9B, 0xDC)
TEAL = RGBColor(0x00, 0xB4, 0xA0)
ORANGE = RGBColor(0xFF, 0x8C, 0x42)
PURPLE = RGBColor(0x9B, 0x59, 0xB6)

W = Inches(13.333)
H = Inches(7.5)


def bg(s, c):
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = c


def box(s, l, t, w, h, c):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = c
    sh.line.fill.background()
    return sh


def rbox(s, l, t, w, h, c):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = c
    sh.line.fill.background()
    return sh


def tx(s, l, t, w, h, text, sz=20, bold=False, color=WHITE, align=PP_ALIGN.LEFT):
    tb = s.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = "Arial"
    p.alignment = align
    p.space_before = Pt(0)
    p.space_after = Pt(0)
    return tb


def tx_multi(s, l, t, w, h, lines_data, align=PP_ALIGN.CENTER):
    """lines_data = [(text, size, bold, color), ...]"""
    tb = s.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    for i, ld in enumerate(lines_data):
        text, sz, bold, color = ld
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(sz)
        p.font.bold = bold
        p.font.color.rgb = color
        p.font.name = "Arial"
        p.alignment = align
        p.space_before = Pt(2)
        p.space_after = Pt(2)
    return tb


def big_num(s, l, t, w, num, label, num_color=GOLD, label_color=GRAY):
    tx(s, l, t, w, Inches(1), num, 72, True, num_color, PP_ALIGN.CENTER)
    tx(s, l, t + Inches(1.0), w, Inches(0.5), label, 22, False, label_color, PP_ALIGN.CENTER)


def card_stat(s, l, t, w, h, num, label, c):
    rbox(s, l, t, w, h, c)
    tx(s, l, t + Inches(0.15), w, Inches(1), num, 56, True, WHITE, PP_ALIGN.CENTER)
    # Split label by \n into separate lines
    label_lines = label.split("\n")
    lines_data = [(line, 18, True, WHITE) for line in label_lines]
    tx_multi(s, l, t + h - Inches(0.3) - Inches(len(label_lines) * 0.25), w, Inches(len(label_lines) * 0.3), lines_data)


def gold_line(s, t):
    box(s, Inches(0), t, W, Inches(0.06), GOLD)


def generate():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    # ═══════════════════════════════════
    # S1: COVER — แบบ premium dark
    # ═══════════════════════════════════
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, BG)
    gold_line(s, Inches(0))

    rbox(s, Inches(4.2), Inches(0.8), Inches(4.9), Inches(0.6), GOLD)
    tx(s, Inches(4.2), Inches(0.82), Inches(4.9), Inches(0.55),
       "HIGH PERFORMANCE HEALTHCARE", 22, True, BG, PP_ALIGN.CENTER)

    tx(s, Inches(0.5), Inches(2.0), Inches(12.3), Inches(1.5),
       "Hospital Claim AI", 80, True, WHITE, PP_ALIGN.CENTER)

    tx(s, Inches(1.5), Inches(3.6), Inches(10.3), Inches(0.6),
       "AI ตรวจ Claim อัตโนมัติ — ลด Deny เพิ่มรายได้ ครบทุกกองทุน", 26, False, GOLD, PP_ALIGN.CENTER)

    # 3 กองทุน badges
    funds = [("ประกันสังคม", BLUE), ("สปสช./UCEP", TEAL), ("ประกันเอกชน", PURPLE)]
    for i, (name, c) in enumerate(funds):
        x = Inches(3.0) + Inches(i * 2.7)
        rbox(s, x, Inches(4.5), Inches(2.3), Inches(0.5), c)
        tx(s, x, Inches(4.52), Inches(2.3), Inches(0.45), name, 18, True, WHITE, PP_ALIGN.CENTER)

    tx(s, Inches(1), Inches(5.5), Inches(11.3), Inches(0.4),
       "โรงพยาบาลพญาไทศรีราชา  |  BDMS Network", 22, False, GRAY, PP_ALIGN.CENTER)
    tx(s, Inches(1), Inches(6.0), Inches(11.3), Inches(0.4),
       "แผนกศัลยกรรมหัวใจและทรวงอก  |  นพ.เกษมสันต์ เกษมวงศ์", 20, False, GRAY, PP_ALIGN.CENTER)

    gold_line(s, Inches(6.8))
    box(s, Inches(0), Inches(6.9), W, Inches(0.6), CARD)
    tx(s, Inches(0), Inches(6.95), W, Inches(0.5),
       "The Ppitch Awards 2026  |  G5 Passion to Innovate", 18, True, GOLD, PP_ALIGN.CENTER)

    # ═══════════════════════════════════
    # S2: ปัญหา — ตัวเลขเจ็บๆ
    # ═══════════════════════════════════
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, BG)

    tx(s, Inches(0.5), Inches(0.3), Inches(12), Inches(0.6),
       "ปัญหา", 44, True, RED)
    tx(s, Inches(0.5), Inches(0.85), Inches(12), Inches(0.5),
       "ทุกเดือน รพ. สูญเสียรายได้จาก Claim ถูก Deny", 24, False, GRAY)

    # 4 big pain stat cards
    pains = [
        ("~10%", "Deny Rate เฉลี่ย", RED),
        ("4.8M", "บาท/ปี สูญเสีย", RED),
        ("30 นาที", "ตรวจ Manual 1 เคส", ORANGE),
        ("2 ชม.", "ร่าง Appeal 1 เคส", ORANGE),
    ]
    for i, (num, label, c) in enumerate(pains):
        x = Inches(0.5) + Inches(i * 3.2)
        card_stat(s, x, Inches(1.8), Inches(2.9), Inches(2.5), num, label, c)

    # Root causes
    tx(s, Inches(0.8), Inches(4.6), Inches(12), Inches(0.5),
       "สาเหตุหลัก:", 26, True, GOLD)

    causes = [
        "ICD Coding ไม่สอดคล้องกับหัตถการ",
        "เอกสารอุปกรณ์ไม่ครบ (ADP/GPO/VMI)",
        "ตรวจ Manual เปิดหลายระบบ ผิดพลาดง่าย",
        "อุทธรณ์ไม่ทัน 15 วัน ร่างหนังสือนาน",
    ]
    for i, c in enumerate(causes):
        col = i % 2
        row = i // 2
        x = Inches(0.8) + Inches(col * 6.3)
        y = Inches(5.3) + Inches(row * 0.7)
        tx(s, x, y, Inches(6), Inches(0.5), f"  {c}", 22, False, WHITE)

    # ═══════════════════════════════════
    # S3: Solution — AI ตรวจ 8 จุด
    # ═══════════════════════════════════
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, LIGHT_BG)

    box(s, Inches(0), Inches(0), W, Inches(1.3), BG)
    tx(s, Inches(0.8), Inches(0.15), Inches(11), Inches(0.6),
       "Solution — AI ตรวจ 8 จุดอัตโนมัติ", 44, True, WHITE)
    tx(s, Inches(0.8), Inches(0.75), Inches(11), Inches(0.4),
       "จาก 30 นาที เหลือไม่ถึง 5 นาที ครอบคลุม 100% ของ checkpoints", 22, False, GOLD)

    checks = [
        ("01", "ข้อมูลพื้นฐาน", "PDx + CID + วันที่", BLUE),
        ("02", "Dx-Proc Match", "ICD-10 ↔ ICD-9-CM", BLUE),
        ("03", "เอกสารอุปกรณ์", "Stent + GPO/VMI", BLUE),
        ("04", "FDH 16-File", "Format + LOS", BLUE),
        ("05", "Timing + Authen", "30 วัน + Authen Code", TEAL),
        ("06", "CC/MCC Optimize", "เพิ่ม DRG Weight", TEAL),
        ("07", "DRG Verify", "Group + RW ถูกต้อง", TEAL),
        ("08", "Drug/Lab Catalog", "TMT + Drug Catalog", TEAL),
    ]
    for i, (num, name, desc, c) in enumerate(checks):
        col = i % 4
        row = i // 4
        x = Inches(0.4) + Inches(col * 3.2)
        y = Inches(1.6) + Inches(row * 2.85)

        rbox(s, x, y, Inches(2.95), Inches(2.5), WHITE)

        # Number circle
        circle = s.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.95), y + Inches(0.15), Inches(0.85), Inches(0.85))
        circle.fill.solid()
        circle.fill.fore_color.rgb = c
        circle.line.fill.background()
        p = circle.text_frame.paragraphs[0]
        p.text = num
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        tx(s, x + Inches(0.1), y + Inches(1.15), Inches(2.75), Inches(0.5),
           name, 24, True, DARK, PP_ALIGN.CENTER)
        tx(s, x + Inches(0.1), y + Inches(1.65), Inches(2.75), Inches(0.5),
           desc, 18, False, GRAY, PP_ALIGN.CENTER)

    # ═══════════════════════════════════
    # S4: 6 คำสั่ง /claim
    # ═══════════════════════════════════
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, BG)

    tx(s, Inches(0.5), Inches(0.3), Inches(12), Inches(0.6),
       "6 คำสั่ง — พิมพ์ครั้งเดียว AI จัดการ", 44, True, WHITE)
    tx(s, Inches(0.5), Inches(0.9), Inches(12), Inches(0.4),
       "เจ้าหน้าที่ใช้งานง่าย ไม่ต้องเรียนนาน", 22, False, GRAY)

    feats = [
        ("/claim check", "ตรวจก่อนส่ง", "Score 0-100%", BLUE),
        ("/claim predict", "ทำนาย Deny", "10 Risk Factors", TEAL),
        ("/claim code", "Smart Coder", "Notes → ICD อัตโนมัติ", GREEN),
        ("/claim deny", "วิเคราะห์ Deny", "Root Cause + วิธีแก้", RED),
        ("/claim appeal", "ร่าง Appeal", "DOCX พร้อมพิมพ์ลงนาม", ORANGE),
        ("/claim batch", "ตรวจทั้งล็อต", "CSV → เรียงตามมูลค่า", PURPLE),
    ]
    for i, (cmd, title, desc, c) in enumerate(feats):
        col = i % 3
        row = i // 3
        x = Inches(0.4) + Inches(col * 4.25)
        y = Inches(1.6) + Inches(row * 2.85)

        rbox(s, x, y, Inches(3.85), Inches(2.5), CARD)

        # Command badge
        rbox(s, x + Inches(0.2), y + Inches(0.2), Inches(2.8), Inches(0.55), c)
        tx(s, x + Inches(0.2), y + Inches(0.22), Inches(2.8), Inches(0.5),
           cmd, 22, True, WHITE, PP_ALIGN.CENTER)

        tx(s, x + Inches(0.3), y + Inches(0.95), Inches(3.3), Inches(0.5),
           title, 28, True, WHITE)
        tx(s, x + Inches(0.3), y + Inches(1.5), Inches(3.3), Inches(0.5),
           desc, 20, False, GRAY)

    # ═══════════════════════════════════
    # S5: Before / After
    # ═══════════════════════════════════
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, LIGHT_BG)

    box(s, Inches(0), Inches(0), W, Inches(1.3), BG)
    tx(s, Inches(0.8), Inches(0.15), Inches(11), Inches(0.6),
       "ก่อน vs หลัง — ผลต่างชัดเจน", 44, True, WHITE)

    # Before
    box(s, Inches(0.5), Inches(1.6), Inches(6), Inches(5.4), WHITE)
    rbox(s, Inches(1.8), Inches(1.75), Inches(3.3), Inches(0.6), RED)
    tx(s, Inches(1.8), Inches(1.78), Inches(3.3), Inches(0.55),
       "BEFORE", 26, True, WHITE, PP_ALIGN.CENTER)

    befores = [
        ("30 นาที", "ตรวจ 1 เคส"),
        ("~10%", "Deny Rate"),
        ("2 ชม.", "ร่าง Appeal"),
        ("-4.8M", "สูญเสีย/ปี"),
    ]
    for i, (num, label) in enumerate(befores):
        y = Inches(2.6) + Inches(i * 1.15)
        tx(s, Inches(0.8), y, Inches(2.8), Inches(0.6), num, 40, True, RED, PP_ALIGN.RIGHT)
        tx(s, Inches(3.8), y + Inches(0.05), Inches(2.5), Inches(0.5), label, 24, False, DARK)

    # After
    box(s, Inches(6.8), Inches(1.6), Inches(6), Inches(5.4), WHITE)
    rbox(s, Inches(8.1), Inches(1.75), Inches(3.3), Inches(0.6), GREEN)
    tx(s, Inches(8.1), Inches(1.78), Inches(3.3), Inches(0.55),
       "AFTER — AI", 26, True, WHITE, PP_ALIGN.CENTER)

    afters = [
        ("5 นาที", "ตรวจ 1 เคส (-83%)"),
        ("<3%", "Deny Rate (target)"),
        ("5 นาที", "Appeal + DOCX (-96%)"),
        ("+2.5M", "รายได้เพิ่ม/ปี"),
    ]
    for i, (num, label) in enumerate(afters):
        y = Inches(2.6) + Inches(i * 1.15)
        tx(s, Inches(7.1), y, Inches(2.8), Inches(0.6), num, 40, True, GREEN, PP_ALIGN.RIGHT)
        tx(s, Inches(10.1), y + Inches(0.05), Inches(2.5), Inches(0.5), label, 24, False, DARK)

    # ═══════════════════════════════════
    # S6: เคสจริง
    # ═══════════════════════════════════
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, BG)

    tx(s, Inches(0.5), Inches(0.3), Inches(12), Inches(0.6),
       "เคสจริง — STEMI Anterior + PCI", 44, True, WHITE)

    # Patient info
    rbox(s, Inches(0.5), Inches(1.2), Inches(6.2), Inches(4.3), CARD)
    tx(s, Inches(0.8), Inches(1.3), Inches(5.6), Inches(0.5),
       "AN 69-03556", 28, True, GOLD)

    info = [
        "PDx: I21.0 Acute MI anterior wall",
        "Proc: 36.06 + 37.22 + 88.56",
        "DRG: 05290 | RW: 8.6544",
        "Charge: 71,322 บาท",
        "สิทธิ: UC (UCEP ฉุกเฉิน)",
    ]
    for i, line in enumerate(info):
        tx(s, Inches(0.8), Inches(2.0) + Inches(i * 0.6), Inches(5.6), Inches(0.5),
           line, 22, False, WHITE)

    rbox(s, Inches(0.8), Inches(5.0), Inches(5.6), Inches(0.5), RED)
    tx(s, Inches(0.8), Inches(5.02), Inches(5.6), Inches(0.45),
       "Deny: HC09 + HC13 + Drug Code ไม่ตรง", 20, True, WHITE, PP_ALIGN.CENTER)

    # AI result
    rbox(s, Inches(7.0), Inches(1.2), Inches(5.8), Inches(4.3), CARD)
    tx(s, Inches(7.3), Inches(1.3), Inches(5.2), Inches(0.5),
       "AI วิเคราะห์ + แนะนำ:", 28, True, GREEN)

    fixes = [
        "แก้ ADP: TYPE=5, Serial ตรง GPO",
        "แก้ DRU: map TMT code ใหม่",
        "เพิ่ม CC: E11.65 (DM) → +8% RW",
        "สร้าง DOCX Appeal พร้อมส่ง",
    ]
    for i, fix in enumerate(fixes):
        tx(s, Inches(7.3), Inches(2.0) + Inches(i * 0.6), Inches(5.2), Inches(0.5),
           fix, 22, False, WHITE)

    # Before/After score
    card_stat(s, Inches(7.2), Inches(4.5), Inches(2.6), Inches(2.0), "65%", "Deny Risk ก่อน", RED)
    card_stat(s, Inches(10.1), Inches(4.5), Inches(2.6), Inches(2.0), "12%", "Deny Risk หลัง", GREEN)

    # ═══════════════════════════════════
    # S7: ผลลัพธ์ — ตัวเลขใหญ่จัด (45% scoring weight)
    # ═══════════════════════════════════
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, BG)

    tx(s, Inches(0.5), Inches(0.2), Inches(12), Inches(0.6),
       "ผลลัพธ์ที่วัดได้", 44, True, GOLD, PP_ALIGN.CENTER)

    # 4 big results
    results = [
        (">90%", "อัตราเบิกผ่านเพิ่มขึ้น", GREEN),
        ("+15%", "DRG Weight เพิ่ม", TEAL),
        ("-83%", "ลดเวลาตรวจ Claim", BLUE),
        ("2.5M+", "บาท/ปี รายได้เพิ่ม", GOLD),
    ]
    for i, (num, label, c) in enumerate(results):
        x = Inches(0.5) + Inches(i * 3.2)
        card_stat(s, x, Inches(1.2), Inches(2.9), Inches(2.6), num, label, c)

    # Detail table
    gold_line(s, Inches(4.1))
    rows = [
        ("ตรวจ Claim", "30 นาที → 5 นาที", "ลด 83%"),
        ("ร่าง Appeal", "2 ชม. → 5 นาที", "ลด 96%"),
        ("ครอบคลุม", "9 แผนก + 3 กองทุน", "20 AI Skills"),
        ("ระบบ Knowledge", "25+ ไฟล์ อัพเดตได้", "Thai DRG v6.3"),
    ]
    for i, (item, detail, result) in enumerate(rows):
        y = Inches(4.5) + Inches(i * 0.7)
        tx(s, Inches(0.8), y, Inches(3.8), Inches(0.5), item, 22, True, WHITE)
        tx(s, Inches(4.8), y, Inches(4.5), Inches(0.5), detail, 22, False, GRAY)
        tx(s, Inches(10.0), y, Inches(2.8), Inches(0.5), result, 22, True, GOLD, PP_ALIGN.RIGHT)

    # ═══════════════════════════════════
    # S8: ความร่วมมือ + ขยายผล (30% scoring weight)
    # ═══════════════════════════════════
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, LIGHT_BG)

    box(s, Inches(0), Inches(0), W, Inches(1.3), BG)
    tx(s, Inches(0.8), Inches(0.15), Inches(11), Inches(0.6),
       "ความร่วมมือ & แผนขยายผล", 44, True, WHITE)
    tx(s, Inches(0.8), Inches(0.75), Inches(11), Inches(0.4),
       "Co-creation ข้ามทีม + Roadmap 4 Phases", 22, False, GOLD)

    # Co-creation
    rbox(s, Inches(0.4), Inches(1.6), Inches(6.1), Inches(2.3), WHITE)
    tx(s, Inches(0.7), Inches(1.7), Inches(5.5), Inches(0.5),
       "ร่วมพัฒนากับ:", 24, True, TEAL)
    collabs = [
        ("Cath Lab Team", "ข้อมูล clinical + validate ผล"),
        ("IT / HIS (SSB)", "Infrastructure + data pipeline"),
        ("UR / Coder", "ทดสอบ + feedback ทุกสัปดาห์"),
    ]
    for i, (team, role) in enumerate(collabs):
        y = Inches(2.3) + Inches(i * 0.5)
        tx(s, Inches(0.7), y, Inches(2.2), Inches(0.4), team, 20, True, DARK)
        tx(s, Inches(3.0), y, Inches(3.3), Inches(0.4), role, 20, False, GRAY)

    # 3 กองทุน
    rbox(s, Inches(6.8), Inches(1.6), Inches(6.1), Inches(2.3), WHITE)
    tx(s, Inches(7.1), Inches(1.7), Inches(5.5), Inches(0.5),
       "รองรับ 3 กองทุน (DRG logic เดียวกัน):", 24, True, TEAL)
    fund_list = [
        ("ประกันสังคม (SSO)", "เคสหลักของ รพ.เอกชน"),
        ("สปสช. / UCEP", "เคสฉุกเฉิน บัตรทอง"),
        ("ประกันเอกชน", "ขยาย Phase 3"),
    ]
    for i, (fund, note) in enumerate(fund_list):
        y = Inches(2.3) + Inches(i * 0.5)
        tx(s, Inches(7.1), y, Inches(2.8), Inches(0.4), fund, 20, True, DARK)
        tx(s, Inches(10.0), y, Inches(2.7), Inches(0.4), note, 20, False, GRAY)

    # Roadmap
    phases = [
        ("Phase 1", "Q3-Q4/69", "ขยายทุก 9 แผนก\nภายใน รพ.", BLUE),
        ("Phase 2", "ปี 70", "+ ประกันสังคม\nBDMS 3-5 รพ.", TEAL),
        ("Phase 3", "ปี 71", "SaaS Platform\n+ ประกันเอกชน", ORANGE),
        ("Phase 4", "ปี 72-73", "รพ. นอกเครือ\nNLP Thai auto", GOLD),
    ]
    for i, (phase, time, desc, c) in enumerate(phases):
        x = Inches(0.4) + Inches(i * 3.25)
        rbox(s, x, Inches(4.3), Inches(3.0), Inches(2.8), c)
        tx(s, x, Inches(4.4), Inches(3.0), Inches(0.5), phase, 24, True, WHITE, PP_ALIGN.CENTER)
        tx(s, x, Inches(4.85), Inches(3.0), Inches(0.4), time, 20, False, WHITE, PP_ALIGN.CENTER)
        tx(s, x, Inches(5.4), Inches(3.0), Inches(1.2), desc, 20, False, WHITE, PP_ALIGN.CENTER)

    # ═══════════════════════════════════
    # S9: AI COMMAND CENTER — Vision
    # ═══════════════════════════════════
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, BG)

    tx(s, Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.7),
       "Vision — AI Command Center", 48, True, GOLD, PP_ALIGN.CENTER)
    tx(s, Inches(0.5), Inches(0.9), Inches(12.3), Inches(0.4),
       "Claim AI คือ Module แรก ของสมองกลาง รพ.", 24, False, WHITE, PP_ALIGN.CENTER)

    modules = [
        ("Claim AI", "ลด Deny\nเพิ่มรายได้", GREEN, "DONE"),
        ("Bed AI", "จัดเตียง\nลด Wait Time", BLUE, "NEXT"),
        ("Staff AI", "จัดเวร\nลด OT Cost", TEAL, "PLAN"),
        ("Supply AI", "คลังยา/วัสดุ\nลด Stock-out", ORANGE, "PLAN"),
        ("Finance AI", "งบประมาณ\nCashflow", GOLD, "PLAN"),
        ("Quality AI", "JCI/HA\nKPI Monitor", PURPLE, "PLAN"),
    ]
    for i, (name, desc, c, status) in enumerate(modules):
        col = i % 3
        row = i // 3
        x = Inches(0.6) + Inches(col * 4.2)
        y = Inches(1.6) + Inches(row * 2.4)

        rbox(s, x, y, Inches(3.7), Inches(2.0), c)

        # Status
        sc = GREEN if status == "DONE" else (BLUE if status == "NEXT" else GRAY)
        rbox(s, x + Inches(2.5), y + Inches(0.1), Inches(1.0), Inches(0.35), sc)
        tx(s, x + Inches(2.5), y + Inches(0.12), Inches(1.0), Inches(0.3),
           status, 16, True, WHITE, PP_ALIGN.CENTER)

        tx(s, Inches(0.2) + x, y + Inches(0.15), Inches(2.2), Inches(0.5),
           name, 28, True, WHITE)
        tx(s, Inches(0.2) + x, y + Inches(0.6), Inches(3.3), Inches(0.8),
           desc, 20, False, WHITE)

    # CEO Dashboard
    gold_line(s, Inches(6.3))
    rbox(s, Inches(3.2), Inches(6.5), Inches(6.9), Inches(0.7), GOLD)
    tx(s, Inches(3.2), Inches(6.53), Inches(6.9), Inches(0.65),
       "CEO Dashboard — เห็นทุกอย่าง ตัดสินใจทันที", 26, True, BG, PP_ALIGN.CENTER)

    # ═══════════════════════════════════
    # S10: CLOSING
    # ═══════════════════════════════════
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, BG)
    gold_line(s, Inches(0))

    tx(s, Inches(0.5), Inches(0.8), Inches(12.3), Inches(1.5),
       "Hospital Claim AI", 80, True, WHITE, PP_ALIGN.CENTER)

    tx(s, Inches(1), Inches(2.5), Inches(11.3), Inches(0.6),
       "ลด Deny เพิ่มรายได้ ครบทุกกองทุน ทุกแผนก", 28, False, GOLD, PP_ALIGN.CENTER)

    # 4 closing stats
    closing = [
        ("20", "AI Skills", TEAL),
        ("3", "กองทุน", BLUE),
        ("9", "แผนก", GREEN),
        ("<5m", "ต่อเคส", ORANGE),
    ]
    for i, (num, label, c) in enumerate(closing):
        x = Inches(0.6) + Inches(i * 3.2)
        card_stat(s, x, Inches(3.3), Inches(2.8), Inches(1.8), num, label, c)

    tx(s, Inches(0.5), Inches(5.5), Inches(12.3), Inches(0.5),
       "1 แผนก → ทุกกองทุน → ทั้ง BDMS → ทั้งประเทศ", 28, True, GOLD, PP_ALIGN.CENTER)

    tx(s, Inches(1), Inches(6.2), Inches(11.3), Inches(0.4),
       "นพ.เกษมสันต์ เกษมวงศ์  |  kasemson_kas@phyathai.com  |  083-574-7685", 20, False, GRAY, PP_ALIGN.CENTER)

    gold_line(s, Inches(6.85))
    box(s, Inches(0), Inches(6.9), W, Inches(0.6), CARD)
    tx(s, Inches(0), Inches(6.95), W, Inches(0.5),
       "The Ppitch Awards 2026  |  G5 Passion to Innovate  |  BDMS", 18, True, GOLD, PP_ALIGN.CENTER)

    # Save
    out = "docs/Ppitch_Awards_2026_v2.pptx"
    prs.save(out)
    print(f"สร้างเสร็จ: {out}")
    print(f"ขนาด: {os.path.getsize(out):,} bytes | Slides: {len(prs.slides)}")


if __name__ == "__main__":
    generate()
